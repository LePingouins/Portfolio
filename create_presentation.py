from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy
INDIGO     = RGBColor(0x1F, 0x3A, 0x8F)   # indigo mid
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)   # bright blue accent
ACCENT2    = RGBColor(0x06, 0xB6, 0xD4)   # cyan accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xFA, 0xFF)
DARK_TEXT  = RGBColor(0x1E, 0x29, 0x3B)
MID_TEXT   = RGBColor(0x55, 0x65, 0x81)
LIGHT_TEXT = RGBColor(0x94, 0xA3, 0xB8)
BAR_BG     = RGBColor(0xE2, 0xE8, 0xF0)
DARK_CARD  = RGBColor(0x16, 0x24, 0x4E)   # slightly lighter navy for cards

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ────────────────────────────────────────────────────────────────

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill, line_color=None, line_pt=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def send_to_back(slide, shape):
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)

def navy_bg(slide):
    bg = rect(slide, 0, 0, W, H, NAVY)
    send_to_back(slide, bg)
    bar = rect(slide, 0, 0, Inches(0.06), H, ACCENT)
    send_to_back(slide, bar)

def light_card(slide, l, t, w, h):
    return rect(slide, l, t, w, h, OFF_WHITE, RGBColor(0xCB, 0xD5, 0xE1), 0.5)

def dark_card(slide, l, t, w, h):
    return rect(slide, l, t, w, h, DARK_CARD, RGBColor(0x2D, 0x4A, 0x8F), 0.75)

def tb(slide, text, l, t, w, h, sz, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, font="Calibri", italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(sz)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return box

def divider(slide, l, t, w, color=ACCENT):
    return rect(slide, l, t, w, Pt(1.5), color)

def thin_rule(slide, l, t, w):
    return rect(slide, l, t, w, Pt(0.75), RGBColor(0x2D, 0x4A, 0x8F))

def skill_bar(slide, label, pct, l, t, bar_w=Inches(3.5)):
    tb(slide, label, l, t, bar_w, Inches(0.3), 11, color=DARK_TEXT)
    rect(slide, l, t + Inches(0.3), bar_w, Inches(0.13), BAR_BG)
    rect(slide, l, t + Inches(0.3), int(bar_w * pct / 100), Inches(0.13), ACCENT)
    tb(slide, f"{pct}%", l + bar_w + Inches(0.07), t + Inches(0.2),
       Inches(0.42), Inches(0.24), 9, color=MID_TEXT)

def section_tag(slide, label):
    """Small accent pill in top-right corner of content slides."""
    rect(slide, W - Inches(1.6), Inches(0.2), Inches(1.1), Inches(0.38), INDIGO)
    tb(slide, label, W - Inches(1.58), Inches(0.22), Inches(1.06), Inches(0.36),
       9, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

def content_header(slide, title, tag=None):
    tb(slide, title, Inches(0.5), Inches(0.2), Inches(9), Inches(0.72),
       34, bold=True, color=WHITE, font="Calibri")
    divider(slide, Inches(0.5), Inches(0.96), Inches(2.2))
    if tag:
        section_tag(slide, tag)

def section_divider_slide(num_str, title, subtitle):
    slide = blank_slide()
    navy_bg(slide)
    # Vertical accent strip on left
    rect(slide, Inches(0.5), Inches(2.0), Inches(0.08), Inches(3.5), ACCENT)
    # Ghost number
    tb(slide, num_str, Inches(0.7), Inches(1.0), Inches(7), Inches(4.5),
       170, bold=True, color=RGBColor(0x14, 0x26, 0x55),
       align=PP_ALIGN.LEFT, font="Calibri")
    # Title
    tb(slide, title, Inches(0.72), Inches(2.4), Inches(10), Inches(1.3),
       58, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri")
    # Subtitle
    tb(slide, subtitle, Inches(0.72), Inches(3.82), Inches(8), Inches(0.55),
       18, color=ACCENT2, align=PP_ALIGN.LEFT, italic=True)
    # Horizontal rule under title
    divider(slide, Inches(0.72), Inches(3.72), Inches(5.0))
    return slide

# ── SLIDE 1 – Title ────────────────────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)

# Right panel — contrasting dark-card strip
rect(slide, Inches(9.5), 0, Inches(3.83), H, DARK_CARD)
rect(slide, Inches(9.5), 0, Pt(1.5), H, ACCENT)

# Badge
badge = rect(slide, Inches(0.5), Inches(1.4), Inches(3.6), Inches(0.4), INDIGO)
tb(slide, "FULL-STACK DEVELOPER  ·  MONTREAL, QC",
   Inches(0.55), Inches(1.43), Inches(3.5), Inches(0.35),
   9.5, bold=True, color=ACCENT2)

# Name — two-line
tb(slide, "Olivier", Inches(0.5), Inches(1.9), Inches(9), Inches(1.25),
   80, bold=True, color=WHITE, font="Calibri")
tb(slide, "Goudreault", Inches(0.5), Inches(3.0), Inches(9), Inches(1.25),
   80, bold=False, color=ACCENT, font="Calibri")

# Subtitle
tb(slide, "Champlain College  ·  Computer Science Technology  ·  2023–2026",
   Inches(0.5), Inches(4.35), Inches(8.8), Inches(0.45),
   13, color=LIGHT_TEXT, italic=True)

# Contact line
tb(slide, "oligoudreault@gmail.com   ·   github.com/LePingouins",
   Inches(0.5), Inches(4.92), Inches(8.8), Inches(0.38),
   12, color=ACCENT2)

# Bottom rule
divider(slide, Inches(0.5), Inches(5.5), Inches(6.5))

# Right panel content
tb(slide, "At a glance", Inches(9.7), Inches(0.6), Inches(3.4), Inches(0.45),
   13, bold=True, color=ACCENT2)
thin_rule(slide, Inches(9.7), Inches(1.08), Inches(3.1))
glance = [
    "🎓 DEC — Computer Science",
    "💻 Full-Stack (Java · React)",
    "🌐 Bilingual FR / EN",
    "🚀 Internship @ Bushido '26",
    "🎮 Gaming & AI enthusiast",
]
for i, line in enumerate(glance):
    tb(slide, line, Inches(9.7), Inches(1.25) + i * Inches(0.72),
       Inches(3.3), Inches(0.6), 12, color=LIGHT_TEXT)

# ── SLIDE 2 – Agenda ─────────────────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)

tb(slide, "What We'll Cover", Inches(0.5), Inches(0.2), W, Inches(0.75),
   34, bold=True, color=WHITE)
divider(slide, Inches(0.5), Inches(0.96), Inches(3.2))

sections = [
    ("01", "About Me",        "Who I am & what drives me"),
    ("02", "Skills",          "Languages, frameworks & tools"),
    ("03", "Projects",        "5 production builds"),
    ("04", "Experience",      "Work history & internship"),
    ("05", "Live Demo & Q&A", "The portfolio website in action"),
]
col_w = Inches(2.28)
gap   = Inches(0.23)
start = Inches(0.5)

for i, (num, head, sub) in enumerate(sections):
    lft = start + i * (col_w + gap)
    dark_card(slide, lft, Inches(1.35), col_w, Inches(4.55))
    rect(slide, lft, Inches(1.35), col_w, Inches(0.1), ACCENT)
    tb(slide, num, lft, Inches(1.5), col_w, Inches(0.95),
       40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    tb(slide, head, lft + Inches(0.12), Inches(2.5), col_w - Inches(0.24), Inches(0.52),
       14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, sub, lft + Inches(0.12), Inches(3.06), col_w - Inches(0.24), Inches(0.8),
       10, color=LIGHT_TEXT, align=PP_ALIGN.CENTER, wrap=True)

# ── SLIDE 3 – Section divider: About ─────────────────────────────────────
section_divider_slide("01", "ABOUT ME", "Who I am & what drives me")

# ── SLIDE 4 – About Me content ───────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)
content_header(slide, "About Me", "01")

# Left column — bio card
light_card(slide, Inches(0.5), Inches(1.15), Inches(7.8), Inches(5.2))

tb(slide, "Hi, I'm Olivier", Inches(0.75), Inches(1.3), Inches(7.3), Inches(0.55),
   21, bold=True, color=DARK_TEXT)
divider(slide, Inches(0.75), Inches(1.87), Inches(2.0), ACCENT)

bio = (
    "I am a passionate software developer with a strong foundation in full-stack "
    "development. I build bilingual, production-grade web applications using Java, "
    "Spring Boot, and React — and I love the challenge of shipping real products "
    "that real people use.\n\n"
    "I recently completed my DEC in Computer Science Technology at Champlain College "
    "(2023–2026). For the next year I will be fully focused on two of my personal "
    "projects. If nothing major comes out of them after a year, I'll be heading to "
    "university to deepen my studies in software engineering.\n\n"
    "Outside of code: gaming, reading, and everything AI & open-source."
)
tb(slide, bio, Inches(0.75), Inches(2.02), Inches(7.3), Inches(3.2),
   11.5, color=MID_TEXT, wrap=True)

# Right column — quick-facts (small accent squares, no ovals)
facts = [
    ("📍", "Montreal, QC, Canada"),
    ("🎓", "Champlain College — DEC Comp. Sci."),
    ("🚀", "Internship @ Bushido (Spring 2026)"),
    ("💡", "Goal: AI & open-source"),
    ("🌐", "Bilingual — French & English"),
]
fy = Inches(1.2)
for icon, fact in facts:
    rect(slide, Inches(8.55), fy + Inches(0.08), Inches(0.06), Inches(0.3), ACCENT)
    tb(slide, icon + "  " + fact, Inches(8.72), fy, Inches(4.35), Inches(0.52),
       11.5, color=WHITE)
    fy += Inches(0.64)

# ── Clean horizontal timeline ─────────────────────────────────────────────
TL_Y = Inches(6.52)
TL_L = Inches(0.5)
TL_W = Inches(12.3)

rect(slide, TL_L, TL_Y, TL_W, Pt(1.2), ACCENT)   # timeline spine

events = [
    (0.0,  "2019",      "Handler\nHorizon Nature"),
    (0.22, "2021",      "Office\nAssistant"),
    (0.50, "2023",      "Champlain\nCollege"),
    (0.78, "2026",      "Bushido\nInternship"),
    (1.0,  "Now",       "Personal\nProjects"),
]
for rel, year, label in events:
    x = TL_L + rel * TL_W
    # dot
    rect(slide, x - Pt(3), TL_Y - Pt(3), Pt(6), Pt(6), WHITE)
    # year above
    tb(slide, year, x - Inches(0.45), TL_Y - Inches(0.4), Inches(0.9), Inches(0.3),
       8.5, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    # label below
    tb(slide, label, x - Inches(0.55), TL_Y + Inches(0.1), Inches(1.1), Inches(0.55),
       7.5, color=LIGHT_TEXT, align=PP_ALIGN.CENTER, wrap=True)

# ── SLIDE 5 – Section divider: Skills ─────────────────────────────────────
section_divider_slide("02", "SKILLS", "The ones that matter most")

# ── SLIDE 6 – Skills content ──────────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)
content_header(slide, "Skills", "02")

COL_Y  = Inches(1.18)
COL_H  = Inches(5.7)
HEAD_H = Inches(0.5)

# ── Languages ──
CX = Inches(0.5)
CW = Inches(3.9)
light_card(slide, CX, COL_Y, CW, COL_H)
rect(slide, CX, COL_Y, CW, Inches(0.07), ACCENT)
tb(slide, "Languages", CX + Inches(0.2), COL_Y + Inches(0.1), CW - Inches(0.3), HEAD_H,
   13, bold=True, color=DARK_TEXT)
langs = [("Java", 90), ("Python", 80), ("SQL", 80), ("TypeScript", 75)]
for j, (name, pct) in enumerate(langs):
    skill_bar(slide, name, pct, CX + Inches(0.2),
              COL_Y + HEAD_H + Inches(0.2) + j * Inches(0.98), bar_w=Inches(3.3))

# ── Frameworks ──
CX = Inches(4.72)
light_card(slide, CX, COL_Y, CW, COL_H)
rect(slide, CX, COL_Y, CW, Inches(0.07), ACCENT2)
tb(slide, "Frameworks", CX + Inches(0.2), COL_Y + Inches(0.1), CW - Inches(0.3), HEAD_H,
   13, bold=True, color=DARK_TEXT)
fws = [("React", 85), ("Spring Boot", 85), ("Next.js", 70)]
for j, (name, pct) in enumerate(fws):
    skill_bar(slide, name, pct, CX + Inches(0.2),
              COL_Y + HEAD_H + Inches(0.2) + j * Inches(0.98), bar_w=Inches(3.3))

# ── Tools ──
CX = Inches(8.93)
light_card(slide, CX, COL_Y, CW, COL_H)
rect(slide, CX, COL_Y, CW, Inches(0.07), RGBColor(0xA7, 0x8B, 0xFA))
tb(slide, "Tools & Platforms", CX + Inches(0.2), COL_Y + Inches(0.1), CW - Inches(0.3), HEAD_H,
   13, bold=True, color=DARK_TEXT)
tools = [("GitHub", 95), ("VS Code", 90), ("IntelliJ", 85), ("Linux", 80)]
for j, (name, pct) in enumerate(tools):
    skill_bar(slide, name, pct, CX + Inches(0.2),
              COL_Y + HEAD_H + Inches(0.2) + j * Inches(0.98), bar_w=Inches(3.3))

# ── SLIDE 7 – Section divider: Projects ──────────────────────────────────
section_divider_slide("03", "PROJECTS", "5 production builds")

# ── SLIDE 8 – Projects grid (3 main) ──────────────────────────────────────
slide = blank_slide()
navy_bg(slide)
content_header(slide, "Projects", "03")

projects_a = [
    {
        "name":  "CourtierPro",
        "tag":   "Real Estate Platform",
        "desc":  "Bilingual broker-client management platform for Nabizada Courtier Inc. "
                 "Production-grade, used by real clients.",
        "tech":  "React · Spring Boot · PostgreSQL",
        "url":   "courtier-pro.ca",
        "color": ACCENT,
    },
    {
        "name":  "Entretien Bâtiment",
        "tag":   "Work Order Management",
        "desc":  "Task & work-order system with km tracking and archive module "
                 "for field-service teams.",
        "tech":  "React · Spring Boot · PostgreSQL",
        "url":   "entretien-batiment.com",
        "color": ACCENT2,
    },
    {
        "name":  "Bushido",
        "tag":   "Dojo Marketing & SaaS",
        "desc":  "Rebuilt the entire marketing website for Bushido Dojo Software and "
                 "implemented project contract pausing and extension features.",
        "tech":  "Ruby on Rails · PostgreSQL",
        "url":   "bushidodojosoftware.com",
        "color": RGBColor(0xA7, 0x8B, 0xFA),
    },
]

CW = Inches(3.87)
GAP = Inches(0.22)
cx  = Inches(0.5)
for p in projects_a:
    light_card(slide, cx, Inches(1.15), CW, Inches(5.85))
    rect(slide, cx, Inches(1.15), CW, Inches(0.1), p["color"])
    # tag
    rect(slide, cx + Inches(0.18), Inches(1.32), Inches(2.35), Inches(0.28), p["color"])
    tb(slide, p["tag"], cx + Inches(0.2), Inches(1.34), Inches(2.31), Inches(0.25),
       8, bold=True, color=WHITE)
    tb(slide, p["name"], cx + Inches(0.18), Inches(1.68), CW - Inches(0.35), Inches(0.55),
       16, bold=True, color=DARK_TEXT)
    tb(slide, p["desc"], cx + Inches(0.18), Inches(2.3), CW - Inches(0.35), Inches(1.85),
       10.5, color=MID_TEXT, wrap=True)
    thin_rule(slide, cx + Inches(0.18), Inches(4.22), CW - Inches(0.4))
    tb(slide, p["tech"], cx + Inches(0.18), Inches(4.35), CW - Inches(0.35), Inches(0.38),
       9.5, color=ACCENT, italic=True)
    tb(slide, "🔗 " + p["url"], cx + Inches(0.18), Inches(4.82), CW - Inches(0.35), Inches(0.38),
       9.5, color=MID_TEXT)
    cx += CW + GAP

# ── SLIDE 9 – Projects grid (2 more) ─────────────────────────────────────
slide = blank_slide()
navy_bg(slide)
content_header(slide, "Projects  (continued)", "03")

projects_b = [
    {
        "name":  "Champlain PetClinic",
        "tag":   "Microservices · Academic",
        "desc":  "Multi-year academic project: Spring Boot microservices architecture "
                 "teaching Scrum, TDD, CI/CD and version control to an entire cohort.",
        "tech":  "React · Spring Boot · TypeScript · Docker",
        "url":   "github.com/cgerard321/champlain_petclinic",
        "color": RGBColor(0x34, 0xD3, 0x99),
    },
    {
        "name":  "Le Pré Paré",
        "tag":   "E-commerce Showcase",
        "desc":  "Product showcase and e-commerce site for a natural-products farm. "
                 "Clean, modern storefront — no chemical fertilizers or pesticides.",
        "tech":  "Next.js · React · TypeScript",
        "url":   "leprepare.ca",
        "color": RGBColor(0xFB, 0xBF, 0x24),
    },
]

BW  = Inches(6.0)
GAP = Inches(0.32)
cx  = Inches(0.5)
for p in projects_b:
    light_card(slide, cx, Inches(1.15), BW, Inches(5.85))
    rect(slide, cx, Inches(1.15), BW, Inches(0.1), p["color"])
    rect(slide, cx + Inches(0.18), Inches(1.32), Inches(2.6), Inches(0.28), p["color"])
    tb(slide, p["tag"], cx + Inches(0.2), Inches(1.34), Inches(2.56), Inches(0.25),
       8, bold=True, color=WHITE)
    tb(slide, p["name"], cx + Inches(0.18), Inches(1.68), BW - Inches(0.35), Inches(0.55),
       18, bold=True, color=DARK_TEXT)
    tb(slide, p["desc"], cx + Inches(0.18), Inches(2.3), BW - Inches(0.35), Inches(2.2),
       11.5, color=MID_TEXT, wrap=True)
    thin_rule(slide, cx + Inches(0.18), Inches(4.6), BW - Inches(0.4))
    tb(slide, p["tech"], cx + Inches(0.18), Inches(4.73), BW - Inches(0.35), Inches(0.38),
       10.5, color=ACCENT, italic=True)
    tb(slide, "🔗 " + p["url"], cx + Inches(0.18), Inches(5.22), BW - Inches(0.35), Inches(0.38),
       10.5, color=MID_TEXT)
    cx += BW + GAP

# ── SLIDE 10 – Section divider: Experience ────────────────────────────────
section_divider_slide("04", "EXPERIENCE", "Work history & internship")

# ── SLIDE 11 – Experience content ────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)
content_header(slide, "Work Experience", "04")

jobs = [
    {
        "period":  "Spring 2026",
        "role":    "Full-Stack Developer Intern",
        "company": "Bushido Dojo Software  ·  Remote",
        "bullets": [
            "Rebuilt the entire marketing website for Bushido Dojo Software.",
            "Implemented project contract pausing and extension features.",
            "Stack: Ruby on Rails · PostgreSQL — live at bushidodojosoftware.com.",
        ],
        "color": ACCENT,
    },
    {
        "period":  "2021 – 2023",
        "role":    "Data Entry Clerk / Office Assistant",
        "company": "Horizon Nature  ·  Saint-Léonard, Montréal, QC",
        "bullets": [
            "Entered and verified high-volume data with precision and confidentiality.",
            "Supported daily operations and coordinated document management.",
        ],
        "color": ACCENT2,
    },
    {
        "period":  "2019 – 2021",
        "role":    "Handler",
        "company": "Horizon Nature  ·  Saint-Léonard, Montréal, QC",
        "bullets": [
            "Movement, organization, and preparation of materials.",
            "Developed reliability, teamwork, and time-management skills.",
        ],
        "color": RGBColor(0xA7, 0x8B, 0xFA),
    },
]

y = Inches(1.18)
for job in jobs:
    light_card(slide, Inches(0.5), y, Inches(12.33), Inches(1.65))
    rect(slide, Inches(0.5), y, Inches(0.1), Inches(1.65), job["color"])
    tb(slide, job["period"], Inches(0.73), y + Inches(0.1), Inches(2.1), Inches(0.35),
       10, bold=True, color=job["color"])
    tb(slide, job["role"], Inches(0.73), y + Inches(0.43), Inches(6.8), Inches(0.42),
       14, bold=True, color=DARK_TEXT)
    tb(slide, job["company"], Inches(0.73), y + Inches(0.84), Inches(6.8), Inches(0.34),
       10, color=MID_TEXT, italic=True)
    bullet_text = "  •  " + "\n  •  ".join(job["bullets"])
    tb(slide, bullet_text, Inches(7.7), y + Inches(0.1), Inches(4.85), Inches(1.5),
       9.5, color=MID_TEXT, wrap=True)
    y += Inches(1.76)

# Education footer bar
rect(slide, Inches(0.5), Inches(6.42), Inches(12.33), Inches(0.75), DARK_CARD)
rect(slide, Inches(0.5), Inches(6.42), Pt(1.5), Inches(0.75), ACCENT2)
tb(slide, "🎓  Champlain College — DEC in Computer Science Technology (2023–2026)   ·   "
         "Next: personal projects · open to new opportunities",
   Inches(0.72), Inches(6.5), Inches(12.0), Inches(0.55),
   10.5, color=WHITE)

# ── SLIDE 12 – Contact ────────────────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)
content_header(slide, "Let's Connect")

# Right panel
rect(slide, Inches(8.8), 0, Inches(4.53), H, DARK_CARD)
rect(slide, Inches(8.8), 0, Pt(1.5), H, ACCENT)

# Contact card (left)
light_card(slide, Inches(0.5), Inches(1.15), Inches(8.1), Inches(5.55))

tb(slide, "Olivier Goudreault", Inches(0.78), Inches(1.32), Inches(7.5), Inches(0.62),
   24, bold=True, color=DARK_TEXT)
divider(slide, Inches(0.78), Inches(1.96), Inches(1.8))
tb(slide, "I'm always open to new opportunities — freelance, full-time, or just a "
          "conversation about tech. Feel free to reach out!",
   Inches(0.78), Inches(2.1), Inches(7.3), Inches(0.85),
   12, color=MID_TEXT, wrap=True)

contacts = [
    ("📧", "EMAIL",    "oligoudreault@gmail.com"),
    ("💼", "LINKEDIN", "linkedin.com/in/olivier-goudreault-09120a386"),
    ("🐙", "GITHUB",   "github.com/LePingouins"),
    ("🌐", "WEBSITE",  "portfolio-frontend-iy8v.onrender.com"),
]
cy = Inches(3.15)
for icon, label, value in contacts:
    rect(slide, Inches(0.78), cy + Inches(0.1), Inches(0.05), Inches(0.35), ACCENT)
    tb(slide, icon + "  " + label, Inches(0.95), cy, Inches(1.8), Inches(0.3),
       8, bold=True, color=ACCENT)
    tb(slide, value, Inches(0.95), cy + Inches(0.3), Inches(7.0), Inches(0.32),
       11.5, color=DARK_TEXT)
    cy += Inches(0.75)

# Right panel — availability box
tb(slide, "AVAILABILITY", Inches(9.05), Inches(0.55), Inches(4.0), Inches(0.4),
   10, bold=True, color=ACCENT2)
thin_rule(slide, Inches(9.05), Inches(0.97), Inches(3.7))
avail_items = [
    "✔  Open to full-time roles",
    "✔  Freelance / contract work",
    "✔  Collaboration on open-source",
    "✔  Coffee chats about tech",
]
for i, item in enumerate(avail_items):
    tb(slide, item, Inches(9.05), Inches(1.15) + i * Inches(0.65),
       Inches(3.85), Inches(0.55), 12, color=LIGHT_TEXT)

# ── SLIDE 13 – Live Demo ──────────────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)

# Top accent band
rect(slide, 0, 0, W, Inches(0.5), DARK_CARD)
rect(slide, 0, 0, W, Pt(2), ACCENT)

tb(slide, "Live Demo", Inches(0), Inches(1.8), W, Inches(1.2),
   66, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
divider(slide, Inches(5.1), Inches(3.1), Inches(3.13))
tb(slide, "portfolio-frontend-iy8v.onrender.com",
   Inches(0), Inches(3.3), W, Inches(0.58),
   19, color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)
tb(slide, "Let's open it together ↓",
   Inches(0), Inches(4.0), W, Inches(0.48),
   13, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)

# ── SLIDE 14 – Thank You ──────────────────────────────────────────────────
slide = blank_slide()
navy_bg(slide)

# Two vertical accent bars for visual interest — no circles
rect(slide, Inches(1.5), Inches(0.4), Pt(2), Inches(6.7), RGBColor(0x1A, 0x32, 0x70))
rect(slide, Inches(11.83), Inches(0.4), Pt(2), Inches(6.7), RGBColor(0x1A, 0x32, 0x70))
rect(slide, Inches(1.5), Inches(0.4), Inches(10.33), Pt(2), RGBColor(0x1A, 0x32, 0x70))
rect(slide, Inches(1.5), Inches(7.1), Inches(10.33), Pt(2), RGBColor(0x1A, 0x32, 0x70))

tb(slide, "Thank You!", Inches(0), Inches(2.1), W, Inches(1.4),
   72, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
divider(slide, Inches(5.1), Inches(3.6), Inches(3.13))
tb(slide, "Any questions?",
   Inches(0), Inches(3.8), W, Inches(0.68),
   28, color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)
tb(slide, "oligoudreault@gmail.com",
   Inches(0), Inches(4.6), W, Inches(0.48),
   15, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────
output_path = r"C:\PortfolioWebsite\Portfolio_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
